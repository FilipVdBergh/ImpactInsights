from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

class powerpointplot:
    def __init__(self):
        self.screen_width = Cm(21)
        self.screen_height = Cm(15)
        self.lane_padding_percentage = 0.1
        self.title_font = "calibri"
        self.title_font_color = (255, 255, 255)
        self.title_font_size = 14
        self.lane_font = "cambria"
        self.lane_font_color = (255, 255, 255)
        self.lane_font_size = 11
        self.lane_color = (33, 66, 99)
        self.lane_alpha = 200
        self.box_title_font = "calibri"
        self.box_title_font_color = (255, 255, 255)

        self.box_title_font_size = 16
        self.box_color = (66, 99, 121)
        self.box_text_font = "calibri"
        self.box_text_font_color = (255, 255, 255)
        self.box_text_font_size = 12
        self.box_width_padding_percentage = 0.10
        self.box_height_padding_percentage = 0.20
        self.bg_color = (255, 255, 255)
        self.bg_picture = None
        self.filename = 'output.pptx'

        self.objects = list()

        self.presentation = Presentation()
        blank_slide_layout = self.presentation.slide_layouts[6]
        self.slide = self.presentation.slides.add_slide(blank_slide_layout)


    def plot(self, canvas):
        self.lane_width = self.screen_width
        self.lane_height = self.screen_height / (len(canvas.layer) + (len(canvas.layer) - 1) * self.lane_padding_percentage)
        maximum_number_of_boxes_per_layer = 0

        if self.bg_picture is not None:
            self.slide.shapes.add_picture(self.bg_picture,0,0, self.screen_width,self.screen_height)


        for l in canvas.layer.values():
            maximum_number_of_boxes_per_layer = max(maximum_number_of_boxes_per_layer, len(l))

        for y, l in enumerate(canvas.layer.values()):
            lane_y = y * (self.lane_height * (1 + self.lane_padding_percentage))
            lane = self.slide.shapes.add_textbox(0, lane_y, self.lane_width, self.lane_height)
            lane.fill.solid()
            lane.fill.fore_color.rgb = RGBColor(*self.lane_color)
            lane.fill.transparency = self.lane_alpha/255

            lane.text = canvas.layernames[y]
            lane.text_frame.paragraphs[0].font.name = self.lane_font
            lane.text_frame.paragraphs[0].font.size = Pt(self.lane_font_size)
            lane.text_frame.paragraphs[0].font.color.rgb = RGBColor(*self.lane_font_color)

            for x, b in enumerate(l):
                box_width = (self.lane_width) / (maximum_number_of_boxes_per_layer + (self.box_width_padding_percentage * maximum_number_of_boxes_per_layer) + self.box_width_padding_percentage)
                box_height =  self.lane_height * (1 - self.box_height_padding_percentage)
                box_x = (self.lane_width - len(l) * box_width - (len(
                    l) + 1) * self.box_width_padding_percentage * box_width) / 2 + self.box_width_padding_percentage * box_width + x * box_width * (
                    1 + self.box_width_padding_percentage)
                box_y = lane_y + self.lane_height * (self.box_height_padding_percentage / 2)

                box = self.slide.shapes.add_textbox(box_x, box_y, box_width, box_height)
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(self.box_color[0],self.box_color[1],self.box_color[2])
                box.text = b.title
                box.text_frame.word_wrap = True
                box.text_frame.paragraphs[0].font.name = self.box_text_font
                box.text_frame.paragraphs[0].font.size = Pt(self.box_text_font_size)
                box.text_frame.paragraphs[0].font.color.rgb = RGBColor(*self.box_text_font_color)


        self.presentation.save(self.filename)
